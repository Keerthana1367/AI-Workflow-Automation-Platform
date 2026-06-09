import PyPDF2
import pandas as pd
from pptx import Presentation
import json

def parse_file(uploaded_file):
    file_type = uploaded_file.name.split(".")[-1].lower()
    text = ""

    if file_type == "pdf":
        reader = PyPDF2.PdfReader(uploaded_file)
        for page in reader.pages:
            text += page.extract_text() or ""
            
    elif file_type in ["txt", "py"]:
        text = uploaded_file.read().decode("utf-8")

    elif file_type in ["csv"]:
        df = pd.read_csv(uploaded_file, nrows=100) # Limit to 100 rows
        text = df.to_string()

    elif file_type in ["xlsx"]:
        df = pd.read_excel(uploaded_file, nrows=100) # Limit to 100 rows
        text = df.to_string()

    elif file_type in ["json"]:
        data = json.load(uploaded_file)
        text = json.dumps(data, indent=2)

    elif file_type in ["pptx"]:
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    elif file_type in ["jpg", "jpeg", "png"]:
        text = "Image parsing is disabled (OCR not supported in this environment)"

    else:
        text = "Unsupported file type"

    # Hard cap at 15,000 characters (~3500 tokens) to prevent Groq Free Tier TPM Limit crashes
    return text[:15000]